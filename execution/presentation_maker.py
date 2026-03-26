#!/usr/bin/env python3
"""
AI Audit — Presentation Maker
Generates a professional 15-slide PPTX from structured JSON data.
Design replicates the established AI Audit presentation template.

Usage:
    python presentation_maker.py --data presentation_data.json --output "AI Audit — Client — Apresentação.pptx"
"""

import argparse
import json
import sys
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =============================================================================
# DESIGN SYSTEM CONSTANTS
# =============================================================================

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Color palette
DARK_NAVY = RGBColor(0x0F, 0x1B, 0x2D)
DARKEST_NAVY = RGBColor(0x0A, 0x16, 0x28)
MEDIUM_DARK = RGBColor(0x1A, 0x30, 0x50)
SLATE_DARK = RGBColor(0x1E, 0x29, 0x3B)
CARD_DARK = RGBColor(0x12, 0x20, 0x33)

SKY_BLUE = RGBColor(0x0E, 0xA5, 0xE9)
DARK_BLUE_ACCENT = RGBColor(0x03, 0x69, 0xA1)

RED = RGBColor(0xEF, 0x44, 0x44)
LIGHT_RED = RGBColor(0xFF, 0x99, 0x99)
RED_BG = RGBColor(0xFF, 0xF0, 0xF0)
DARK_RED_TEXT = RGBColor(0x92, 0x40, 0x0E)

AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE7)

GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
GREEN_LIGHT_BG = RGBColor(0xF0, 0xFD, 0xF4)
DARK_GREEN = RGBColor(0x04, 0x78, 0x57)

SLATE_GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_SLATE = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_TINT = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_LIGHT_TINT = RGBColor(0xF0, 0xF9, 0xFF)

# Semantic color map
COLOR_MAP = {
    "green": GREEN,
    "red": RED,
    "amber": AMBER,
    "blue": SKY_BLUE,
}

BG_COLOR_MAP = {
    "green": GREEN_LIGHT_BG,
    "red": RED_BG,
    "amber": AMBER_BG,
    "blue": BLUE_TINT,
}

FONT_NAME = "Calibri"

# Standard positions (in inches)
LEFT_MARGIN = 0.50
TOP_BAR_HEIGHT = 0.10
SECTION_LABEL_TOP = 0.22
TITLE_TOP = 0.60
CONTENT_TOP = 1.30


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_slide_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color: Optional[RGBColor] = None,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a shape with optional fill color."""
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    # Minimal corner rounding
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color: Optional[RGBColor] = None):
    """Add a rectangle (no rounded corners)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text: str,
                 font_size: int = 17, font_color: RGBColor = SLATE_DARK,
                 bold: bool = False, italic: bool = False,
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap: bool = True):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    if anchor:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic

    try:
        tf.word_wrap = word_wrap
        from pptx.oxml.ns import qn
        bodyPr = tf.paragraphs[0]._p.getparent()
        if bodyPr.tag.endswith('txBody'):
            bp = bodyPr.find(qn('a:bodyPr'))
            if bp is not None:
                anchor_map = {
                    MSO_ANCHOR.TOP: 't',
                    MSO_ANCHOR.MIDDLE: 'ctr',
                    MSO_ANCHOR.BOTTOM: 'b',
                }
                bp.set('anchor', anchor_map.get(anchor, 't'))
    except Exception:
        pass

    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines: list,
                           default_size: int = 17, default_color: RGBColor = SLATE_DARK,
                           alignment=PP_ALIGN.LEFT, line_spacing: float = 1.0):
    """Add a text box with multiple formatted lines.
    Each line in lines is a dict: {text, size, color, bold, italic} (all optional except text).
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            line_data = {"text": line_data}

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = alignment
        p.space_before = Pt(0)
        p.space_after = Pt(2)

        run = p.add_run()
        run.text = line_data["text"]
        run.font.name = FONT_NAME
        run.font.size = Pt(line_data.get("size", default_size))
        run.font.color.rgb = line_data.get("color", default_color)
        run.font.bold = line_data.get("bold", False)
        run.font.italic = line_data.get("italic", False)

    return txBox


def add_top_bar(slide, color: RGBColor):
    """Add thin color bar at top of slide."""
    add_rect(slide, 0, 0, 10.0, TOP_BAR_HEIGHT, color)


def add_left_accent(slide, color: RGBColor = SKY_BLUE):
    """Add left accent bar (for dark slides)."""
    add_rect(slide, 0, 0, 0.14, 5.625, color)


def add_bottom_bar(slide, color: RGBColor = DARKEST_NAVY):
    """Add bottom bar (for dark slides)."""
    add_rect(slide, 0, 5.30, 10.0, 0.325, color)


def add_section_label(slide, text: str, color: RGBColor = SKY_BLUE):
    """Add section label at standard position."""
    add_text_box(slide, LEFT_MARGIN, SECTION_LABEL_TOP, 9.0, 0.35,
                 text, font_size=14, font_color=color, bold=True)


def add_slide_title(slide, text: str, color: RGBColor = SLATE_DARK, font_size: int = 36):
    """Add slide title at standard position."""
    add_text_box(slide, LEFT_MARGIN, TITLE_TOP, 9.0, 0.55,
                 text, font_size=font_size, font_color=color, bold=True)


def add_card(slide, left, top, width, height, accent_color: RGBColor = SKY_BLUE):
    """Add a card with background and accent strip at top."""
    # Card background
    add_shape(slide, left, top, width, height, CARD_BG)
    # Accent strip
    add_rect(slide, left, top, width, 0.07, accent_color)


def add_card_dark(slide, left, top, width, height, accent_color: Optional[RGBColor] = None):
    """Add a card for dark slides."""
    add_shape(slide, left, top, width, height, CARD_DARK)
    if accent_color:
        add_rect(slide, left, top, width, 0.07, accent_color)


def get_semantic_color(color_name: str) -> RGBColor:
    """Get RGBColor from semantic name."""
    return COLOR_MAP.get(color_name, SKY_BLUE)


def get_semantic_bg(color_name: str) -> RGBColor:
    """Get background tint from semantic name."""
    return BG_COLOR_MAP.get(color_name, BLUE_TINT)


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_slide_01_cover(prs, data: Dict):
    """Slide 1: Title/Cover (dark background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, DARK_NAVY)
    add_left_accent(slide, SKY_BLUE)
    add_bottom_bar(slide)

    # "AI Audit" - large title
    add_text_box(slide, 0.80, 1.20, 8.5, 1.20,
                 "AI Audit", font_size=92, font_color=WHITE, bold=True)

    # Client name
    add_text_box(slide, 0.80, 2.50, 8.5, 0.80,
                 data["client_name"], font_size=50, font_color=SKY_BLUE, bold=False)

    # Subtitle line
    add_text_box(slide, 0.80, 3.30, 8.5, 0.40,
                 "Diagnóstico & Oportunidades de Automação com IA",
                 font_size=18, font_color=SLATE_GRAY)

    # Footer
    add_text_box(slide, 0.80, 5.35, 4.0, 0.25,
                 "Confidencial", font_size=12, font_color=LIGHT_SLATE)


def build_slide_02_agenda(prs, data: Dict):
    """Slide 2: Agenda (dark background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_left_accent(slide, SKY_BLUE)
    add_bottom_bar(slide)

    # Title
    add_text_box(slide, 0.80, 0.50, 8.5, 0.60,
                 "Próximos 30 minutos", font_size=42, font_color=WHITE, bold=True)

    agenda = data.get("agenda", {})
    blocks = [agenda.get(f"block_{i}", {}) for i in range(1, 5)]

    y_start = 1.40
    block_height = 0.85
    gap = 0.15

    for i, block in enumerate(blocks):
        y = y_start + i * (block_height + gap)

        # Time badge
        add_shape(slide, 0.80, y, 1.20, 0.70, MEDIUM_DARK)
        add_text_box(slide, 0.80, y + 0.10, 1.20, 0.50,
                     block.get("time", f"{i*8}-{(i+1)*8} min"),
                     font_size=15, font_color=SKY_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, 2.20, y + 0.05, 6.5, 0.35,
                     block.get("title", ""),
                     font_size=24, font_color=WHITE, bold=True)

        # Description
        add_text_box(slide, 2.20, y + 0.38, 6.5, 0.30,
                     block.get("description", ""),
                     font_size=15, font_color=LIGHT_SLATE)


def build_slide_03_diagnostico(prs, data: Dict):
    """Slide 3: Diagnostic KPIs (light background, 3 cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, SKY_BLUE)

    add_section_label(slide, "DIAGNÓSTICO", SKY_BLUE)
    add_slide_title(slide, "Panorama Atual")

    kpis = data.get("diagnostico_kpis", [])

    card_width = 2.90
    gap = 0.15
    start_x = 0.40
    card_top = 1.50
    card_height = 3.40

    for i, kpi in enumerate(kpis[:3]):
        x = start_x + i * (card_width + gap)
        accent = get_semantic_color(kpi.get("color", "blue"))

        add_card(slide, x, card_top, card_width, card_height, accent)

        # Value (big number)
        add_text_box(slide, x + 0.20, card_top + 0.30, card_width - 0.40, 0.80,
                     kpi.get("value", "—"),
                     font_size=50, font_color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, x + 0.20, card_top + 1.20, card_width - 0.40, 0.40,
                     kpi.get("label", ""),
                     font_size=22, font_color=SLATE_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + 0.20, card_top + 1.70, card_width - 0.40, 1.40,
                     kpi.get("description", ""),
                     font_size=15, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)


def build_slide_04_descobertas_resumo(prs, data: Dict):
    """Slide 4: 3 Critical Discoveries summary (light background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, RED)

    add_section_label(slide, "ANÁLISE DIAGNÓSTICA", RED)
    add_slide_title(slide, "3 Descobertas Críticas")

    discoveries = data.get("descobertas_criticas", [])

    card_width = 2.90
    gap = 0.15
    start_x = 0.40
    card_top = 1.50
    card_height = 3.50

    for i, disc in enumerate(discoveries[:3]):
        x = start_x + i * (card_width + gap)
        accent = get_semantic_color(disc.get("color", "red"))

        add_card(slide, x, card_top, card_width, card_height, accent)

        # Number label
        add_text_box(slide, x + 0.20, card_top + 0.20, card_width - 0.40, 0.30,
                     disc.get("number", f"0{i+1}"),
                     font_size=14, font_color=LIGHT_RED, bold=True)

        # Title
        add_text_box(slide, x + 0.20, card_top + 0.55, card_width - 0.40, 0.40,
                     disc.get("title", ""),
                     font_size=22, font_color=SLATE_DARK, bold=True)

        # Subtitle
        add_text_box(slide, x + 0.20, card_top + 1.00, card_width - 0.40, 0.35,
                     disc.get("subtitle", ""),
                     font_size=15, font_color=SLATE_GRAY)

        # Bullets
        bullets = disc.get("bullets", [])
        bullet_text = "\n".join([f"• {b}" for b in bullets])
        add_text_box(slide, x + 0.20, card_top + 1.50, card_width - 0.40, 1.70,
                     bullet_text,
                     font_size=14, font_color=SLATE_GRAY)


def build_slide_05_descoberta_1(prs, data: Dict):
    """Slide 5: Discovery 1 Detail (light background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, AMBER)

    detail = data.get("descoberta_1_detail", {})

    add_section_label(slide, detail.get("section_label", "DESCOBERTA 01"), AMBER)
    add_slide_title(slide, detail.get("title", "Descoberta 1"), font_size=36)

    # Metric cards (up to 4 in a row)
    metrics = detail.get("metrics", [])
    card_width = 2.18
    gap = 0.12
    start_x = 0.40
    card_top = 1.40

    for i, metric in enumerate(metrics[:4]):
        x = start_x + i * (card_width + gap)
        add_card(slide, x, card_top, card_width, 1.10, AMBER)

        add_text_box(slide, x + 0.15, card_top + 0.15, card_width - 0.30, 0.50,
                     metric.get("value", "—"),
                     font_size=28, font_color=SLATE_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        label_text = metric.get("label", "")
        if metric.get("sublabel"):
            label_text += f"\n{metric['sublabel']}"
        add_text_box(slide, x + 0.15, card_top + 0.60, card_width - 0.30, 0.45,
                     label_text,
                     font_size=12, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Quote box
    quote = detail.get("quote", {})
    if quote:
        quote_top = 2.80
        # Vertical accent bar
        add_rect(slide, 0.40, quote_top, 0.06, 1.10, AMBER)

        add_text_box(slide, 0.60, quote_top + 0.05, 5.50, 0.70,
                     f'"{quote.get("text", "")}"',
                     font_size=15, font_color=SLATE_DARK, italic=True)

        add_text_box(slide, 0.60, quote_top + 0.75, 5.50, 0.30,
                     f'— {quote.get("attribution", "")}',
                     font_size=13, font_color=SLATE_GRAY)

    # Callout box
    callout = detail.get("callout", "")
    if callout:
        callout_top = 4.10
        add_shape(slide, 0.40, callout_top, 9.10, 0.80, BLUE_TINT)
        add_text_box(slide, 0.60, callout_top + 0.10, 8.70, 0.60,
                     f"◆ {callout}",
                     font_size=14, font_color=DARK_BLUE_ACCENT)


def build_slide_06_descoberta_2(prs, data: Dict):
    """Slide 6: Discovery 2 Detail (light background) - process flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, AMBER)

    detail = data.get("descoberta_2_detail", {})

    add_section_label(slide, detail.get("section_label", "DESCOBERTA 02"), AMBER)
    add_slide_title(slide, detail.get("title", "Descoberta 2"), font_size=36)

    # Flow steps (horizontal process flow)
    steps = detail.get("flow_steps", [])
    step_width = 1.80
    arrow_width = 0.40
    start_x = 0.40
    step_top = 1.50
    step_height = 0.90

    for i, step in enumerate(steps[:4]):
        x = start_x + i * (step_width + arrow_width)

        # Step badge (number)
        add_shape(slide, x, step_top, 0.40, 0.40, MEDIUM_DARK)
        add_text_box(slide, x, step_top + 0.02, 0.40, 0.36,
                     step.get("step", str(i + 1)),
                     font_size=18, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Step label
        add_text_box(slide, x + 0.50, step_top + 0.05, step_width - 0.55, 0.35,
                     step.get("label", ""),
                     font_size=17, font_color=SLATE_DARK, bold=True)

        # Arrow between steps
        if i < len(steps) - 1:
            arrow_x = x + step_width + 0.05
            add_text_box(slide, arrow_x, step_top + 0.05, 0.30, 0.35,
                         "→", font_size=25, font_color=LIGHT_SLATE,
                         alignment=PP_ALIGN.CENTER)

    # Storage locations
    locations = detail.get("storage_locations", [])
    if locations:
        loc_top = 2.70
        add_text_box(slide, 0.40, loc_top - 0.30, 9.0, 0.30,
                     "Informação Fragmentada em:",
                     font_size=14, font_color=SLATE_GRAY, bold=True)

        loc_width = 2.90
        for i, loc in enumerate(locations[:3]):
            x = 0.40 + i * (loc_width + 0.15)
            add_shape(slide, x, loc_top, loc_width, 0.60, AMBER_BG)
            add_text_box(slide, x + 0.15, loc_top + 0.10, loc_width - 0.30, 0.40,
                         loc, font_size=15, font_color=DARK_RED_TEXT,
                         alignment=PP_ALIGN.CENTER)

    # Bottom stats
    stats = detail.get("bottom_stats", [])
    stat_top = 3.70
    stat_width = 2.90
    stat_height = 1.20

    for i, stat in enumerate(stats[:3]):
        x = 0.40 + i * (stat_width + 0.15)
        add_card(slide, x, stat_top, stat_width, stat_height, AMBER)

        add_text_box(slide, x + 0.15, stat_top + 0.20, stat_width - 0.30, 0.50,
                     stat.get("value", "—"),
                     font_size=28, font_color=SLATE_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, stat_top + 0.70, stat_width - 0.30, 0.40,
                     stat.get("label", ""),
                     font_size=14, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)


def build_slide_07_descoberta_3(prs, data: Dict):
    """Slide 7: Discovery 3 Detail (light background) - funnel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, AMBER)

    detail = data.get("descoberta_3_detail", {})

    add_section_label(slide, detail.get("section_label", "DESCOBERTA 03"), AMBER)
    add_slide_title(slide, detail.get("title", "Descoberta 3"), font_size=36)

    # Funnel visualization (left side)
    funnel = detail.get("funnel_stages", [])
    funnel_left = 0.50
    funnel_top = 1.60
    max_width = 4.50
    stage_height = 0.70
    gap = 0.10

    for i, stage in enumerate(funnel[:5]):
        width_pct = stage.get("width_pct", 100 - i * 20) / 100
        w = max_width * width_pct
        offset = (max_width - w) / 2
        y = funnel_top + i * (stage_height + gap)

        # Funnel bar (darker blue gradient)
        bar_colors = [SKY_BLUE, RGBColor(0x38, 0xBD, 0xF8), AMBER, RGBColor(0xFB, 0xBF, 0x24), RED]
        bar_color = bar_colors[i] if i < len(bar_colors) else SLATE_GRAY

        add_shape(slide, funnel_left + offset, y, w, stage_height, bar_color)

        # Label + value on the bar
        label = stage.get("label", "")
        value = stage.get("value", "")
        add_text_box(slide, funnel_left + offset + 0.15, y + 0.10, w - 0.30, 0.25,
                     label, font_size=14, font_color=WHITE, bold=True)
        add_text_box(slide, funnel_left + offset + 0.15, y + 0.35, w - 0.30, 0.25,
                     value, font_size=13, font_color=WHITE)

    # Info panel (right side)
    info = detail.get("info_panel", {})
    if info:
        panel_left = 5.40
        panel_top = 1.60
        panel_width = 4.20
        panel_height = 3.50

        add_shape(slide, panel_left, panel_top, panel_width, panel_height, CARD_BG)
        add_rect(slide, panel_left, panel_top, panel_width, 0.07, SKY_BLUE)

        add_text_box(slide, panel_left + 0.20, panel_top + 0.20, panel_width - 0.40, 0.35,
                     info.get("title", ""),
                     font_size=20, font_color=SLATE_DARK, bold=True)

        bullets = info.get("bullets", [])
        bullet_text = "\n\n".join([f"• {b}" for b in bullets])
        add_text_box(slide, panel_left + 0.20, panel_top + 0.65, panel_width - 0.40, 2.60,
                     bullet_text,
                     font_size=14, font_color=SLATE_GRAY)


def build_slide_08_custo_status_quo(prs, data: Dict):
    """Slide 8: Cost of Status Quo (dark background, 3 big numbers)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_top_bar(slide, RED)

    add_section_label(slide, "IMPACTO", RED)
    add_slide_title(slide, "O Custo do Status Quo", WHITE, font_size=42)

    costs = data.get("custo_status_quo", [])
    card_width = 2.90
    gap = 0.15
    start_x = 0.40
    card_top = 1.70
    card_height = 3.10

    for i, cost in enumerate(costs[:3]):
        x = start_x + i * (card_width + gap)
        add_card_dark(slide, x, card_top, card_width, card_height, RED)

        add_text_box(slide, x + 0.20, card_top + 0.50, card_width - 0.40, 0.80,
                     cost.get("value", "—"),
                     font_size=50, font_color=RED, bold=True,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + 0.20, card_top + 1.40, card_width - 0.40, 0.40,
                     cost.get("label", ""),
                     font_size=20, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + 0.20, card_top + 1.90, card_width - 0.40, 0.90,
                     cost.get("sublabel", ""),
                     font_size=15, font_color=LIGHT_SLATE,
                     alignment=PP_ALIGN.CENTER)


def build_slide_09_benchmarks(prs, data: Dict):
    """Slide 9: Competitive Benchmarks (light background, table grid)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, SKY_BLUE)

    add_section_label(slide, "CONTEXTO COMPETITIVO", SKY_BLUE)
    add_slide_title(slide, "Benchmarks do Mercado")

    benchmarks = data.get("benchmarks", {})
    metrics = benchmarks.get("metrics", [])
    client_vals = benchmarks.get("client_values", [])
    avg_vals = benchmarks.get("market_avg", [])
    top_vals = benchmarks.get("top_25", [])
    gaps = benchmarks.get("gap", [])

    # Column headers
    headers = ["Métrica", data.get("client_name", "Cliente"), "Média Mercado", "Top 25%", "Gap"]
    col_widths = [2.20, 1.70, 1.70, 1.70, 1.70]
    start_x = 0.30
    table_top = 1.40
    header_height = 0.55
    row_height = 0.70

    # Header row
    x = start_x
    for j, (header, cw) in enumerate(zip(headers, col_widths)):
        add_shape(slide, x, table_top, cw - 0.05, header_height, SLATE_DARK)
        add_text_box(slide, x + 0.10, table_top + 0.10, cw - 0.25, header_height - 0.20,
                     header, font_size=14, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        x += cw

    # Data rows
    num_rows = min(len(metrics), 5)
    for i in range(num_rows):
        y = table_top + header_height + 0.08 + i * (row_height + 0.08)
        row_data = [
            metrics[i] if i < len(metrics) else "",
            client_vals[i] if i < len(client_vals) else "",
            avg_vals[i] if i < len(avg_vals) else "",
            top_vals[i] if i < len(top_vals) else "",
            gaps[i] if i < len(gaps) else "",
        ]

        x = start_x
        for j, (val, cw) in enumerate(zip(row_data, col_widths)):
            bg = CARD_BG
            text_color = SLATE_DARK

            # Color the gap column
            if j == 4:
                if val.startswith("-") or val.startswith("−"):
                    bg = RED_BG
                    text_color = RED
                elif val.startswith("+"):
                    bg = GREEN_LIGHT_BG
                    text_color = GREEN

            add_shape(slide, x, y, cw - 0.05, row_height, bg)
            add_text_box(slide, x + 0.10, y + 0.15, cw - 0.25, row_height - 0.30,
                         val, font_size=15, font_color=text_color, bold=(j == 0),
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            x += cw


def build_slide_10_prioridades(prs, data: Dict):
    """Slide 10: 5 Priorities Overview (light background, 5 columns)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, GREEN)

    add_section_label(slide, "ROADMAP", GREEN)
    add_slide_title(slide, "As 5 Prioridades")

    priorities = data.get("prioridades", [])
    col_width = 1.78
    gap = 0.10
    start_x = 0.30
    col_top = 1.50
    col_height = 3.60

    for i, prio in enumerate(priorities[:5]):
        x = start_x + i * (col_width + gap)
        accent = get_semantic_color(prio.get("color", "blue"))

        # Column background
        add_shape(slide, x, col_top, col_width, col_height, CARD_BG)
        add_rect(slide, x, col_top, col_width, 0.50, accent)

        # Priority number
        add_text_box(slide, x, col_top + 0.05, col_width, 0.40,
                     prio.get("number", f"P{i+1}"),
                     font_size=25, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Name
        add_text_box(slide, x + 0.10, col_top + 0.65, col_width - 0.20, 0.55,
                     prio.get("name", ""),
                     font_size=14, font_color=SLATE_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # ROI
        add_text_box(slide, x + 0.10, col_top + 1.30, col_width - 0.20, 0.25,
                     "ROI", font_size=12, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.10, col_top + 1.50, col_width - 0.20, 0.35,
                     prio.get("roi", "—"),
                     font_size=21, font_color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Payback
        add_text_box(slide, x + 0.10, col_top + 1.95, col_width - 0.20, 0.25,
                     "Payback", font_size=12, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.10, col_top + 2.15, col_width - 0.20, 0.30,
                     prio.get("payback", "—"),
                     font_size=17, font_color=SLATE_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Timeline
        add_text_box(slide, x + 0.10, col_top + 2.55, col_width - 0.20, 0.25,
                     "Timeline", font_size=12, font_color=SLATE_GRAY,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.10, col_top + 2.75, col_width - 0.20, 0.30,
                     prio.get("timeline", "—"),
                     font_size=15, font_color=SLATE_DARK,
                     alignment=PP_ALIGN.CENTER)

        # Phase badge
        phase = prio.get("phase", "")
        if phase:
            add_shape(slide, x + 0.15, col_top + 3.15, col_width - 0.30, 0.30, get_semantic_bg(prio.get("color", "blue")))
            add_text_box(slide, x + 0.15, col_top + 3.17, col_width - 0.30, 0.26,
                         phase, font_size=11, font_color=accent, bold=True,
                         alignment=PP_ALIGN.CENTER)


def build_slide_11_quick_wins(prs, data: Dict):
    """Slide 11: P1 & P2 Quick Wins Detail (light background, 2 cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, GREEN)

    add_section_label(slide, "QUICK WINS", GREEN)
    add_slide_title(slide, "P1 & P2 — Vitórias Rápidas", font_size=28)

    items = data.get("quick_wins_detail", [])
    card_width = 4.45
    gap = 0.20
    start_x = 0.40
    card_top = 1.30
    card_height = 3.90

    for i, item in enumerate(items[:2]):
        x = start_x + i * (card_width + gap)
        accent = GREEN if i == 0 else SKY_BLUE

        add_card(slide, x, card_top, card_width, card_height, accent)

        # Priority badge
        add_shape(slide, x + 0.15, card_top + 0.20, 0.50, 0.30, accent)
        add_text_box(slide, x + 0.15, card_top + 0.21, 0.50, 0.28,
                     item.get("priority", f"P{i+1}"),
                     font_size=14, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + 0.75, card_top + 0.20, card_width - 0.95, 0.35,
                     item.get("title", ""),
                     font_size=18, font_color=SLATE_DARK, bold=True)

        # Content rows
        rows = [
            ("Problema:", item.get("problem", "")),
            ("Solução:", item.get("solution", "")),
            ("Ferramentas:", item.get("tools", "")),
            ("Investimento:", item.get("investment", "")),
            ("Resultado:", item.get("expected_result", "")),
            ("Timeline:", item.get("timeline", "")),
        ]

        y_offset = card_top + 0.70
        for label, value in rows:
            if value:
                add_text_box(slide, x + 0.20, y_offset, card_width - 0.40, 0.22,
                             label, font_size=11, font_color=accent, bold=True)
                add_text_box(slide, x + 0.20, y_offset + 0.20, card_width - 0.40, 0.35,
                             value, font_size=13, font_color=SLATE_GRAY)
                y_offset += 0.52


def build_slide_12_transformation(prs, data: Dict):
    """Slide 12: P3, P4 & P5 Transformation Detail (light background, 3 cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, SKY_BLUE)

    add_section_label(slide, "TRANSFORMAÇÃO", SKY_BLUE)
    add_slide_title(slide, "P3, P4 & P5 — Transformação Estrutural", font_size=28)

    items = data.get("transformation_detail", [])
    card_width = 3.00
    gap = 0.12
    start_x = 0.30
    card_top = 1.30
    card_height = 3.90

    colors = [SKY_BLUE, AMBER, AMBER]

    for i, item in enumerate(items[:3]):
        x = start_x + i * (card_width + gap)
        accent = colors[i] if i < len(colors) else SKY_BLUE

        add_card(slide, x, card_top, card_width, card_height, accent)

        # Priority badge
        add_shape(slide, x + 0.12, card_top + 0.18, 0.45, 0.28, accent)
        add_text_box(slide, x + 0.12, card_top + 0.19, 0.45, 0.26,
                     item.get("priority", f"P{i+3}"),
                     font_size=13, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + 0.65, card_top + 0.18, card_width - 0.80, 0.32,
                     item.get("title", ""),
                     font_size=16, font_color=SLATE_DARK, bold=True)

        # Content rows
        rows = [
            ("Problema:", item.get("problem", "")),
            ("Solução:", item.get("solution", "")),
            ("Ferramentas:", item.get("tools", "")),
            ("Investimento:", item.get("investment", "")),
            ("Resultado:", item.get("expected_result", "")),
            ("Timeline:", item.get("timeline", "")),
        ]

        y_offset = card_top + 0.62
        for label, value in rows:
            if value:
                add_text_box(slide, x + 0.15, y_offset, card_width - 0.30, 0.20,
                             label, font_size=10, font_color=accent, bold=True)
                add_text_box(slide, x + 0.15, y_offset + 0.18, card_width - 0.30, 0.32,
                             value, font_size=12, font_color=SLATE_GRAY)
                y_offset += 0.48


def build_slide_13_roi(prs, data: Dict):
    """Slide 13: Consolidated ROI (dark background, 4 summary cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_top_bar(slide, GREEN)

    add_section_label(slide, "BUSINESS CASE", GREEN)
    add_slide_title(slide, "ROI Consolidado — 12 meses", WHITE, font_size=36)

    roi = data.get("roi_consolidado", {})
    items = [
        roi.get("investment", {"value": "—", "label": "Investimento Total"}),
        roi.get("return", {"value": "—", "label": "Retorno Estimado"}),
        roi.get("roi", {"value": "—", "label": "ROI"}),
        roi.get("payback", {"value": "—", "label": "Payback Médio"}),
    ]

    card_width = 2.18
    gap = 0.12
    start_x = 0.40
    card_top = 1.80
    card_height = 2.80

    colors = [SKY_BLUE, GREEN, GREEN, AMBER]

    for i, item in enumerate(items[:4]):
        x = start_x + i * (card_width + gap)
        accent = colors[i]

        add_card_dark(slide, x, card_top, card_width, card_height, accent)

        add_text_box(slide, x + 0.15, card_top + 0.50, card_width - 0.30, 0.80,
                     item.get("value", "—"),
                     font_size=36, font_color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + 0.15, card_top + 1.40, card_width - 0.30, 0.40,
                     item.get("label", ""),
                     font_size=17, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)


def build_slide_14_resultados(prs, data: Dict):
    """Slide 14: Expected Results Table (light background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, GREEN)

    add_section_label(slide, "PROJEÇÃO", GREEN)
    add_slide_title(slide, "Resultados Esperados — 12 Meses")

    results = data.get("resultados_esperados", [])

    # Table headers
    headers = ["Métrica", "Antes", "Depois", "Variação"]
    col_widths = [3.00, 2.20, 2.20, 2.00]
    start_x = 0.30
    table_top = 1.40
    header_height = 0.50
    row_height = 0.60

    # Header row
    x = start_x
    for header, cw in zip(headers, col_widths):
        add_shape(slide, x, table_top, cw - 0.05, header_height, SLATE_DARK)
        add_text_box(slide, x + 0.10, table_top + 0.08, cw - 0.25, header_height - 0.16,
                     header, font_size=14, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        x += cw

    # Data rows
    for i, row in enumerate(results[:7]):
        y = table_top + header_height + 0.08 + i * (row_height + 0.06)
        row_data = [
            row.get("metric", ""),
            row.get("before", ""),
            row.get("after", ""),
            row.get("variation", ""),
        ]

        x = start_x
        for j, (val, cw) in enumerate(zip(row_data, col_widths)):
            bg = CARD_BG
            text_color = SLATE_DARK

            if j == 3:  # Variation column
                if "+" in val:
                    bg = GREEN_LIGHT_BG
                    text_color = GREEN
                elif "-" in val:
                    bg = RED_BG
                    text_color = RED

            add_shape(slide, x, y, cw - 0.05, row_height, bg)
            add_text_box(slide, x + 0.10, y + 0.12, cw - 0.25, row_height - 0.24,
                         val, font_size=15, font_color=text_color, bold=(j == 0),
                         alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
            x += cw


def build_slide_15_proximos_passos(prs, data: Dict):
    """Slide 15: Next Steps (dark background, 3 action items)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_left_accent(slide, SKY_BLUE)
    add_bottom_bar(slide)

    add_text_box(slide, 0.80, 0.50, 8.5, 0.70,
                 "PRÓXIMOS PASSOS", font_size=45, font_color=WHITE, bold=True)

    steps = data.get("proximos_passos", [])
    step_top = 1.50
    step_height = 1.15
    gap = 0.15

    for i, step in enumerate(steps[:3]):
        y = step_top + i * (step_height + gap)

        # Number circle
        add_shape(slide, 0.80, y, 0.60, 0.60, SKY_BLUE, MSO_SHAPE.OVAL)
        add_text_box(slide, 0.80, y + 0.08, 0.60, 0.44,
                     step.get("number", str(i + 1)),
                     font_size=28, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Timeline badge
        timeline = step.get("timeline", "")
        if timeline:
            add_shape(slide, 1.55, y + 0.05, 1.30, 0.30, MEDIUM_DARK)
            add_text_box(slide, 1.55, y + 0.06, 1.30, 0.28,
                         timeline, font_size=12, font_color=SKY_BLUE, bold=True,
                         alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, 1.55, y + 0.40, 7.5, 0.35,
                     step.get("title", ""),
                     font_size=22, font_color=WHITE, bold=True)

        # Description
        add_text_box(slide, 1.55, y + 0.75, 7.5, 0.35,
                     step.get("description", ""),
                     font_size=15, font_color=LIGHT_SLATE)


# =============================================================================
# MAIN
# =============================================================================

def build_presentation(data: Dict, output_path: str):
    """Build the complete 15-slide presentation from structured data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    builders = [
        build_slide_01_cover,
        build_slide_02_agenda,
        build_slide_03_diagnostico,
        build_slide_04_descobertas_resumo,
        build_slide_05_descoberta_1,
        build_slide_06_descoberta_2,
        build_slide_07_descoberta_3,
        build_slide_08_custo_status_quo,
        build_slide_09_benchmarks,
        build_slide_10_prioridades,
        build_slide_11_quick_wins,
        build_slide_12_transformation,
        build_slide_13_roi,
        build_slide_14_resultados,
        build_slide_15_proximos_passos,
    ]

    for builder in builders:
        builder(prs, data)

    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser(description="AI Audit — Presentation Generator")
    parser.add_argument("--data", required=True, help="Path to JSON data file")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    args = parser.parse_args()

    with open(args.data, "r", encoding="utf-8") as f:
        data = json.load(f)

    build_presentation(data, args.output)


if __name__ == "__main__":
    main()
